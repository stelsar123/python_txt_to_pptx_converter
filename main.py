from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.enum.text import  MSO_AUTO_SIZE, MSO_ANCHOR
import os


#settings
black= RGBColor(0,0,0)
white= RGBColor(250,250,250)
lexicon = {"1":"first","2":"second","3":"third","4":"fourth","5":"fifth"} #Continue this lexicon so you can fit your maximum number of questions in your files
PATH = 'INSERT HERE THE PATH OF THE FILES FOLDER'

#Text alteration function
def fonts(frame, margin_top, margin_left, text, fonts, size , bold, color):

    text_frame = frame.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(margin_top)
    text_frame.margin_left = Inches(margin_left)
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = fonts
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color

#Insert slide function. DISCLAIMER: Adjust the values to your needs
def add_slide(prs, layout, questions, index):
    """Return slide newly added to `prs` using `layout` and having `title`."""
    title_slide_layout = prs.slide_layouts[layout]
    slide = prs.slides.add_slide(title_slide_layout)

    #left, top, width, height
    title = slide.shapes.add_textbox(Inches(2.6), Inches(0.45), Inches(5), Inches(1))
    number = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(1), Inches(1))
    body = slide.shapes.add_textbox(Inches(1.8), Inches(2.5), Inches(7.2), Inches(3))

    notes_slide = slide.notes_slide
    #Adjust the following to your needs
    try:
        fonts(number,0, 0, f"{index + 1}", "Arial", 28, True, white)
        fonts(title, 0, 0 , questions[0], "Arial", 20, True, white)
        text = " "
        testing = questions
        text2 = testing[0] + ".\n"
        testing.pop(0)
        for i in testing:
            text2 += i + "\n"
        

        for i in questions:
            text += i + "\n"
        fonts(body, 0 , 0, text, "Arial", 15, False, black)



        #Notes Text
        comments = notes_slide.notes_text_frame
        comments.text = lexicon[f"{index+1}"]+" Question.\n" + text2 
    except:
        return ValueError
    return slide


def load(file, QandA, template):
    try:
        #Create the list of lines in the txt file
        lista =[]
        with open(PATH + "\\"+ file, "r", encoding="utf_8") as f:
            for line in f:
                a=line.strip('\n')
                lista.append(a)

        #Get the size of the list and locate the indexes with the seperator ----
        size = len(lista)
        idx_list = [idx + 1 for idx, val in
                    enumerate(lista) if val == "----"]
        print(lista)
        #Replace the ---- symbol with the space " " and create sublists of questions and answers
        for i in idx_list :
            lista[i-1] = " "
        for i in range(size):
            lista[i] = lista[i].strip(" ")
        res=[]
        helper = []
        for item in lista:
            if item != "":
                helper.append(item)
            else:
                res.append(helper)
                helper = []
        #create the presentation acording to your template
        prs = Presentation(template)


        #add a differnt slide for each question and answer pair

        for j, item in enumerate(res):
            add_slide(prs, 5, item, j)
            
            
        #delete the template slide
        xml_slides = prs.slides._sldIdLst  
        slid = list(xml_slides)
        xml_slides.remove(slid[0]) 
        #rename the file and save it
        final = list(QandA)
        del final[-9:-5]
        output=""
        for i in final:
            output += i

        prs.save(output)
    except:
        return ValueError



def load_multiple(path,template):
    """Get The list of files in the files folder, and load them one by one"""


    mypath = path #convert the variable to avoid error of matching the .path attribute of os library
    onlyfiles = [f for f in os.listdir(mypath) if os.path.isfile(os.path.join(mypath, f))]

    for file in onlyfiles:
        QandA = file+".pptx"
        load(file,QandA,template)
 


if __name__=="__main__":
#Select the template .pptx file
    templ = input("Insert template: ")
    templ = str(templ)
    load_multiple(PATH, templ )
